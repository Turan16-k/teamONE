"""
T5: Dinamik .pptx Üretim Servisi
Finansal veriler ve AI analizlerini alarak şablona işlenmiş sunum üretir.
"""
import io
from decimal import Decimal
from typing import Optional, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from app.models.financial import FinancialReport
from app.models.company import Company


# Renk paleti
COLOR_PRIMARY = RGBColor(0x1A, 0x3C, 0x5E)    # Koyu lacivert
COLOR_SECONDARY = RGBColor(0x2E, 0x86, 0xAB)   # Mavi
COLOR_ACCENT = RGBColor(0xF0, 0x6C, 0x00)      # Turuncu
COLOR_SUCCESS = RGBColor(0x27, 0xAE, 0x60)     # Yeşil
COLOR_DANGER = RGBColor(0xE7, 0x4C, 0x3C)      # Kırmızı
COLOR_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)       # Açık gri
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _fmt(value: Optional[Any], currency: str = "TRY") -> str:
    if value is None:
        return "N/A"
    try:
        num = float(value)
        if abs(num) >= 1_000_000_000:
            return f"{num/1_000_000_000:.2f}Bn {currency}"
        elif abs(num) >= 1_000_000:
            return f"{num/1_000_000:.2f}M {currency}"
        elif abs(num) >= 1_000:
            return f"{num/1_000:.1f}K {currency}"
        return f"{num:,.2f} {currency}"
    except (TypeError, ValueError):
        return "N/A"


def _pct(value: Optional[Any]) -> str:
    if value is None:
        return "N/A"
    try:
        return f"{float(value) * 100:.1f}%"
    except (TypeError, ValueError):
        return "N/A"


def _set_slide_background(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_box(slide, text: str, top: float, left: float, width: float, height: float,
                    font_size: int = 24, bold: bool = True, color: RGBColor = COLOR_PRIMARY) -> None:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_kpi_box(slide, label: str, value: str, top: float, left: float,
                  bg_color: RGBColor = COLOR_SECONDARY) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(2.2), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.05)
    tf.margin_left = Inches(0.1)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(9)
    r2.font.color.rgb = COLOR_LIGHT


def _add_table(slide, headers: list, rows: list, top: float, left: float, width: float) -> None:
    col_count = len(headers)
    row_count = len(rows) + 1
    table = slide.shapes.add_table(row_count, col_count, Inches(left), Inches(top),
                                    Inches(width), Inches(0.35 * row_count)).table

    col_width = Inches(width / col_count)
    for i in range(col_count):
        table.columns[i].width = col_width

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

    for i, row in enumerate(rows):
        bg = COLOR_LIGHT if i % 2 == 0 else COLOR_WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(9)
            run.font.color.rgb = COLOR_PRIMARY


class PptxService:
    """T5: Finansal rapor ve AI analizini .pptx formatında sunar."""

    def generate_financial_report(
        self,
        company: Company,
        report: FinancialReport,
        ai_analysis: Optional[dict] = None,
        currency: str = "TRY",
    ) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # tamamen boş layout

        self._add_cover_slide(prs, blank_layout, company, report)
        self._add_balance_sheet_slide(prs, blank_layout, report, currency)
        self._add_income_statement_slide(prs, blank_layout, report, currency)
        self._add_cash_flow_slide(prs, blank_layout, report, currency)

        if ai_analysis:
            self._add_ratio_analysis_slide(prs, blank_layout, ai_analysis)
            self._add_assessment_slide(prs, blank_layout, ai_analysis)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _add_cover_slide(self, prs, layout, company: Company, report: FinancialReport) -> None:
        slide = prs.slides.add_slide(layout)
        _set_slide_background(slide, COLOR_PRIMARY)

        _add_title_box(slide, company.name, 1.5, 1.0, 11.0, 1.5,
                        font_size=40, bold=True, color=COLOR_WHITE)
        _add_title_box(slide, "Finansal Analiz Raporu", 3.0, 1.0, 11.0, 0.8,
                        font_size=24, bold=False, color=COLOR_SECONDARY)

        period_label = f"{report.fiscal_year} - {report.period.value.upper()}"
        _add_title_box(slide, period_label, 3.9, 1.0, 6.0, 0.6,
                        font_size=18, bold=False, color=COLOR_LIGHT)

        generated_at = datetime.utcnow().strftime("%d.%m.%Y")
        _add_title_box(slide, f"Oluşturulma Tarihi: {generated_at}", 6.5, 1.0, 6.0, 0.5,
                        font_size=12, bold=False, color=COLOR_LIGHT)

        if report.is_ai_generated:
            _add_title_box(slide, "AI destekli analiz içerir", 6.5, 8.0, 4.0, 0.5,
                            font_size=10, bold=False, color=COLOR_ACCENT)

    def _add_balance_sheet_slide(self, prs, layout, report: FinancialReport, currency: str) -> None:
        slide = prs.slides.add_slide(layout)
        _add_title_box(slide, "Bilanço", 0.2, 0.3, 12.0, 0.7, font_size=28, color=COLOR_PRIMARY)
        _add_title_box(slide, "Varlıklar", 1.1, 0.3, 5.5, 0.4, font_size=14, color=COLOR_SECONDARY)

        asset_rows = [
            ("Nakit ve Eşdeğerleri", _fmt(report.cash_and_equivalents, currency)),
            ("Kısa Vadeli Yatırımlar", _fmt(report.short_term_investments, currency)),
            ("Alacaklar", _fmt(report.accounts_receivable, currency)),
            ("Stoklar", _fmt(report.inventory, currency)),
            ("Diğer Dönen Varlıklar", _fmt(report.other_current_assets, currency)),
            ("TOPLAM DÖNEN VARLIKLAR", _fmt(report.total_current_assets, currency)),
            ("Maddi Duran Varlıklar", _fmt(report.property_plant_equipment, currency)),
            ("Maddi Olmayan Duran Varlıklar", _fmt(report.intangible_assets, currency)),
            ("TOPLAM DURAN VARLIKLAR", _fmt(report.total_non_current_assets, currency)),
            ("TOPLAM VARLIKLAR", _fmt(report.total_assets, currency)),
        ]
        _add_table(slide, ["Kalem", "Tutar"], asset_rows, 1.5, 0.3, 5.5)

        _add_title_box(slide, "Yükümlülükler & Özkaynak", 1.1, 6.3, 6.5, 0.4,
                        font_size=14, color=COLOR_SECONDARY)
        liab_rows = [
            ("Ticari Borçlar", _fmt(report.accounts_payable, currency)),
            ("Kısa Vadeli Borçlar", _fmt(report.short_term_debt, currency)),
            ("TOPLAM KISA VADELİ YÜK.", _fmt(report.total_current_liabilities, currency)),
            ("Uzun Vadeli Borçlar", _fmt(report.long_term_debt, currency)),
            ("TOPLAM UZUN VADELİ YÜK.", _fmt(report.total_non_current_liabilities, currency)),
            ("TOPLAM YÜKÜMLÜLÜKLER", _fmt(report.total_liabilities, currency)),
            ("Ödenmiş Sermaye", _fmt(report.share_capital, currency)),
            ("Geçmiş Yıl Karları", _fmt(report.retained_earnings, currency)),
            ("TOPLAM ÖZKAYNAKLAR", _fmt(report.total_equity, currency)),
        ]
        _add_table(slide, ["Kalem", "Tutar"], liab_rows, 1.5, 6.3, 6.5)

        # KPI kutucukları
        kpis = [
            ("Toplam Varlık", _fmt(report.total_assets, currency), 0.3),
            ("Toplam Borç", _fmt(report.total_liabilities, currency), 2.6),
            ("Özkaynak", _fmt(report.total_equity, currency), 4.9),
        ]
        for label, value, left in kpis:
            _add_kpi_box(slide, label, value, 6.2, left)

    def _add_income_statement_slide(self, prs, layout, report: FinancialReport, currency: str) -> None:
        slide = prs.slides.add_slide(layout)
        _add_title_box(slide, "Gelir Tablosu", 0.2, 0.3, 12.0, 0.7, font_size=28, color=COLOR_PRIMARY)

        rows = [
            ("Gelir (Net Satışlar)", _fmt(report.revenue, currency)),
            ("Satılan Malın Maliyeti", _fmt(report.cost_of_goods_sold, currency)),
            ("Brüt Kar", _fmt(report.gross_profit, currency)),
            ("Faaliyet Giderleri", _fmt(report.operating_expenses, currency)),
            ("EBITDA", _fmt(report.ebitda, currency)),
            ("EBIT (Faaliyet Karı)", _fmt(report.ebit, currency)),
            ("Faiz Giderleri", _fmt(report.interest_expense, currency)),
            ("Vergi Öncesi Kar", _fmt(report.income_before_tax, currency)),
            ("Vergi", _fmt(report.income_tax, currency)),
            ("NET KAR", _fmt(report.net_income, currency)),
        ]
        _add_table(slide, ["Kalem", "Tutar"], rows, 1.1, 0.3, 7.0)

        kpis = [
            ("Gelir", _fmt(report.revenue, currency), 7.8, COLOR_SECONDARY),
            ("Brüt Kar", _fmt(report.gross_profit, currency), 7.8, COLOR_SUCCESS),
            ("Net Kar", _fmt(report.net_income, currency), 7.8, COLOR_ACCENT),
        ]
        tops = [1.1, 2.4, 3.7]
        for (label, value, left, color), top in zip(kpis, tops):
            _add_kpi_box(slide, label, value, top, left, bg_color=color)

    def _add_cash_flow_slide(self, prs, layout, report: FinancialReport, currency: str) -> None:
        slide = prs.slides.add_slide(layout)
        _add_title_box(slide, "Nakit Akış Tablosu", 0.2, 0.3, 12.0, 0.7, font_size=28, color=COLOR_PRIMARY)

        rows = [
            ("Operasyonel Nakit Akışı", _fmt(report.operating_cash_flow, currency)),
            ("Yatırım Nakit Akışı", _fmt(report.investing_cash_flow, currency)),
            ("Finansman Nakit Akışı", _fmt(report.financing_cash_flow, currency)),
            ("Serbest Nakit Akışı", _fmt(report.free_cash_flow, currency)),
            ("Net Nakit Değişimi", _fmt(report.net_change_in_cash, currency)),
        ]
        _add_table(slide, ["Kalem", "Tutar"], rows, 1.1, 0.3, 7.5)

        kpis = [
            ("Operasyonel", _fmt(report.operating_cash_flow, currency), 8.2, COLOR_SECONDARY),
            ("Serbest", _fmt(report.free_cash_flow, currency), 10.5, COLOR_SUCCESS),
        ]
        for label, value, left, color in kpis:
            _add_kpi_box(slide, label, value, 1.1, left, bg_color=color)

    def _add_ratio_analysis_slide(self, prs, layout, ai_analysis: dict) -> None:
        slide = prs.slides.add_slide(layout)
        _add_title_box(slide, "Rasyo Analizi", 0.2, 0.3, 12.0, 0.7, font_size=28, color=COLOR_PRIMARY)

        sections = [
            ("Likidite Rasyoları", "liquidity", [
                ("Cari Oran", "current_ratio"),
                ("Asit-Test Oranı", "quick_ratio"),
                ("Nakit Oranı", "cash_ratio"),
            ], 0.3),
            ("Kaldıraç Rasyoları", "leverage", [
                ("Borç/Özkaynak", "debt_to_equity"),
                ("Borç/Varlık", "debt_to_assets"),
                ("Faiz Karşılama", "interest_coverage"),
            ], 4.5),
            ("Kârlılık Rasyoları", "profitability", [
                ("Brüt Kar Marjı", "gross_margin"),
                ("Net Kar Marjı", "net_margin"),
                ("EBITDA Marjı", "ebitda_margin"),
                ("ROA", "roa"),
                ("ROE", "roe"),
            ], 8.7),
        ]

        for title, key, metrics, left in sections:
            _add_title_box(slide, title, 1.0, left, 4.0, 0.4, font_size=13, color=COLOR_SECONDARY)
            section_data = ai_analysis.get(key, {})
            rows = []
            for label, field in metrics:
                val = section_data.get(field)
                if val is None:
                    rows.append((label, "N/A"))
                elif isinstance(val, (int, float)):
                    # Oran ise float, marj ise yüzde
                    if "margin" in field or "ratio" in field or "roa" in field or "roe" in field:
                        if val < 5:
                            rows.append((label, f"{val:.2f}x"))
                        else:
                            rows.append((label, _pct(val)))
                    else:
                        rows.append((label, f"{val:.2f}"))
                else:
                    rows.append((label, str(val)))
            _add_table(slide, ["Rasyo", "Değer"], rows, 1.5, left, 4.0)

    def _add_assessment_slide(self, prs, layout, ai_analysis: dict) -> None:
        slide = prs.slides.add_slide(layout)
        _add_title_box(slide, "AI Değerlendirme Özeti", 0.2, 0.3, 12.0, 0.7,
                        font_size=28, color=COLOR_PRIMARY)

        overall = ai_analysis.get("overall_assessment", "")
        if overall:
            _add_title_box(slide, "Genel Değerlendirme", 1.1, 0.3, 12.5, 0.4,
                            font_size=14, color=COLOR_SECONDARY)
            _add_title_box(slide, overall, 1.6, 0.3, 12.5, 1.2,
                            font_size=11, bold=False, color=COLOR_PRIMARY)

        risks = ai_analysis.get("risk_indicators", [])
        positives = ai_analysis.get("positive_indicators", [])

        if risks:
            _add_title_box(slide, "Risk Göstergeleri", 3.0, 0.3, 6.0, 0.4,
                            font_size=13, color=COLOR_DANGER)
            risk_text = "\n".join(f"• {r}" for r in risks[:5])
            _add_title_box(slide, risk_text, 3.5, 0.3, 6.0, 2.0,
                            font_size=10, bold=False, color=COLOR_PRIMARY)

        if positives:
            _add_title_box(slide, "Güçlü Yönler", 3.0, 6.8, 6.0, 0.4,
                            font_size=13, color=COLOR_SUCCESS)
            pos_text = "\n".join(f"• {p}" for p in positives[:5])
            _add_title_box(slide, pos_text, 3.5, 6.8, 6.0, 2.0,
                            font_size=10, bold=False, color=COLOR_PRIMARY)


pptx_service = PptxService()
