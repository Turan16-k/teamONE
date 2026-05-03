"""
Tier'a göre çıktı dosyaları üretici.
basic      → JSON
professional → JSON + PDF (grafikler dahil)
enterprise   → JSON + PDF + PPTX + Excel
"""
from __future__ import annotations
import io
import json
import logging
from datetime import datetime
from typing import Dict, Optional

from .models import AnalysisTier, AnalysisResult

logger = logging.getLogger(__name__)


class TieredReporter:
    def generate(
        self,
        result: AnalysisResult,
        charts: Dict[str, bytes],
        raw_analysis: dict,
    ) -> Dict[str, bytes]:
        files: Dict[str, bytes] = {}

        files["analiz_sonucu.json"] = self._json(result, raw_analysis)

        if result.tier in (AnalysisTier.PROFESSIONAL, AnalysisTier.ENTERPRISE):
            pdf = self._pdf(result, charts, result.tier)
            if pdf:
                files["finansal_rapor.pdf"] = pdf

        if result.tier == AnalysisTier.ENTERPRISE:
            pptx = self._pptx(result, charts)
            if pptx:
                files["sunum.pptx"] = pptx
            xlsx = self._excel(result, raw_analysis)
            if xlsx:
                files["veri_tablosu.xlsx"] = xlsx

        for name, data in charts.items():
            if data:
                files[f"grafikler/{name}"] = data

        return files

    # ── JSON ──────────────────────────────────────────────────────────────

    def _json(self, result: AnalysisResult, raw: dict) -> bytes:
        payload = {
            "meta": {
                "analiz_tarihi": datetime.now().isoformat(),
                "paket": result.tier.value,
                "belge_sayisi": result.document_count,
            },
            "ozet": result.executive_summary,
            "ana_bulgular": result.key_insights,
            "belge_ozeti": [s.model_dump() for s in result.document_summaries],
        }
        if result.financial_metrics:
            payload["finansal_metrikler"] = result.financial_metrics.model_dump()
        if result.detailed_analysis:
            payload["kapsamli_analiz"] = result.detailed_analysis
        if result.recommendations:
            payload["oneriler"] = result.recommendations
        if result.anomalies:
            payload["anomaliler"] = result.anomalies
        if raw.get("ratio_analysis"):
            payload["rasyo_analizi"] = raw["ratio_analysis"]
        return json.dumps(payload, ensure_ascii=False, indent=2).encode("utf-8")

    # ── PDF ───────────────────────────────────────────────────────────────

    def _pdf(self, result: AnalysisResult, charts: Dict[str, bytes], tier: AnalysisTier) -> Optional[bytes]:
        try:
            return self._pdf_reportlab(result, charts, tier)
        except Exception as exc:
            logger.warning("reportlab hatası, yedek PDF deneniyor: %s", exc)
            try:
                return self._pdf_fpdf(result)
            except Exception as exc2:
                logger.error("PDF oluşturulamadı: %s", exc2)
                return None

    def _pdf_reportlab(self, result: AnalysisResult, charts: Dict[str, bytes], tier: AnalysisTier) -> bytes:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Image,
            Table, TableStyle, HRFlowable, PageBreak,
        )
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

        W, H = A4
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4,
                                rightMargin=2*cm, leftMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)

        styles = getSampleStyleSheet()
        BLUE = colors.HexColor("#1565C0")

        def S(name, parent="Normal", **kw):
            return ParagraphStyle(name, parent=styles[parent], **kw)

        s_title  = S("T",  "Title",   fontSize=20, textColor=BLUE, spaceAfter=14)
        s_h1     = S("H1", "Heading1", fontSize=14, textColor=BLUE, spaceBefore=14, spaceAfter=7)
        s_h2     = S("H2", "Heading2", fontSize=11, textColor=colors.HexColor("#424242"),
                     spaceBefore=9, spaceAfter=4)
        s_body   = S("B",  "Normal",   fontSize=10, leading=16, alignment=TA_JUSTIFY)
        s_bullet = S("BL", "Normal",   fontSize=10, leftIndent=18, spaceBefore=2)
        s_risk   = S("R",  "Normal",   fontSize=10, leftIndent=18, textColor=colors.HexColor("#C62828"))
        s_ok     = S("OK", "Normal",   fontSize=10, leftIndent=18, textColor=colors.HexColor("#2E7D32"))

        tier_label = {"basic": "Temel", "professional": "Profesyonel", "enterprise": "Kurumsal"}.get(
            tier.value, "Analiz")
        story = [
            Spacer(1, 0.8*cm),
            Paragraph("Finansal Belge Analiz Raporu", s_title),
            Paragraph(
                f"<font size='11' color='#616161'>{tier_label} Analiz &nbsp;|&nbsp; "
                f"{datetime.now().strftime('%d.%m.%Y %H:%M')}</font>",
                styles["Normal"],
            ),
            HRFlowable(width="100%", thickness=2, color=BLUE, spaceAfter=18),
            Paragraph("Yönetici Özeti", s_h1),
            Paragraph(result.executive_summary, s_body),
            Spacer(1, 0.4*cm),
        ]

        if result.key_insights:
            story += [Paragraph("Ana Bulgular", s_h1)]
            for ins in result.key_insights:
                story.append(Paragraph(f"&bull; &nbsp;{ins}", s_bullet))
            story.append(Spacer(1, 0.4*cm))

        # Finansal metrikler tablosu
        if result.financial_metrics:
            m = result.financial_metrics
            c = m.currency

            def fmt(v):
                return f"{v:,.0f} {c}" if v else "-"

            rows = [
                ["Metrik", "Değer"],
                ["Gelir",                   fmt(m.revenue)],
                ["Gider",                   fmt(m.expenses)],
                ["Net Kâr / Zarar",         fmt(m.net_income)],
                ["Toplam Varlık",           fmt(m.total_assets)],
                ["Toplam Borç",             fmt(m.total_liabilities)],
                ["Özkaynak",                fmt(m.total_equity)],
                ["İşletme Nakit Akışı",     fmt(m.operating_cash_flow)],
                ["Dönem",                   m.period or "-"],
            ]
            tbl = Table(rows, colWidths=[8*cm, 8*cm])
            tbl.setStyle(TableStyle([
                ("BACKGROUND",  (0, 0), (-1, 0), BLUE),
                ("TEXTCOLOR",   (0, 0), (-1, 0), colors.white),
                ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F5F5")]),
                ("GRID",        (0, 0), (-1, -1), 0.5, colors.HexColor("#BDBDBD")),
                ("PADDING",     (0, 0), (-1, -1), 7),
                ("ALIGN",       (1, 1), (1, -1), "RIGHT"),
            ]))
            story += [Paragraph("Finansal Metrikler", s_h1), tbl, Spacer(1, 0.4*cm)]

        # Grafikler
        if charts:
            story.append(PageBreak())
            story.append(Paragraph("Grafikler", s_h1))
            for cname, cbytes in charts.items():
                if cbytes:
                    label = cname.replace(".png", "").replace("_", " ").title()
                    story += [
                        Paragraph(label, s_h2),
                        Image(io.BytesIO(cbytes), width=15.5*cm, height=8.5*cm),
                        Spacer(1, 0.4*cm),
                    ]

        # Belge bazı özet
        if result.document_summaries:
            story.append(PageBreak())
            story.append(Paragraph("Belge Bazlı Analiz", s_h1))
            for i, ds in enumerate(result.document_summaries, 1):
                story.append(Paragraph(f"{i}. {ds.filename}", s_h2))
                story.append(Paragraph(ds.brief_summary, s_body))
                for f in ds.key_findings:
                    story.append(Paragraph(f"• {f}", s_bullet))
                story.append(Spacer(1, 0.3*cm))

        if result.detailed_analysis:
            story += [PageBreak(), Paragraph("Kapsamlı Analiz", s_h1),
                      Paragraph(result.detailed_analysis, s_body)]

        if result.recommendations:
            story += [Spacer(1, 0.4*cm), Paragraph("Öneriler", s_h1)]
            for rec in result.recommendations:
                story.append(Paragraph(f"✓ &nbsp;{rec}", s_ok))

        if result.anomalies:
            story += [Spacer(1, 0.4*cm), Paragraph("Risk ve Anomaliler", s_h1)]
            for a in result.anomalies:
                story.append(Paragraph(f"⚠ &nbsp;{a}", s_risk))

        doc.build(story)
        buf.seek(0)
        return buf.read()

    def _pdf_fpdf(self, result: AnalysisResult) -> bytes:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "Finansal Belge Analiz Raporu", ln=True, align="C")
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(0, 7, datetime.now().strftime("%d.%m.%Y %H:%M"), ln=True, align="C")
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Yönetici Ozeti", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 6, result.executive_summary.encode("latin-1", "replace").decode("latin-1"))
        return bytes(pdf.output())

    # ── PPTX ─────────────────────────────────────────────────────────────

    def _pptx(self, result: AnalysisResult, charts: Dict[str, bytes]) -> Optional[bytes]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.warning("python-pptx kurulu değil")
            return None

        BLUE  = RGBColor(0x15, 0x65, 0xC0)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        DARK  = RGBColor(0x21, 0x21, 0x21)
        GREEN = RGBColor(0x2E, 0x7D, 0x32)
        RED   = RGBColor(0xC6, 0x28, 0x28)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def txb(slide, text, l, t, w, h, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf  = box.text_frame; tf.word_wrap = True
            p   = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color

        def bg(slide, hex_str="1565C0"):
            fill = slide.background.fill; fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(hex_str)

        # Kapak
        sl = prs.slides.add_slide(blank); bg(sl, "0D47A1")
        txb(sl, "FİNANSAL BELGE ANALİZ RAPORU", 0.5, 1.8, 12.3, 1.3,
            size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(sl, "Kurumsal Analiz Paketi", 0.5, 3.2, 12.3, 0.8,
            size=16, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)
        txb(sl, datetime.now().strftime("%d %B %Y"), 0.5, 4.1, 12.3, 0.6,
            size=13, color=RGBColor(0xB0, 0xBE, 0xC5), align=PP_ALIGN.CENTER)

        # Yönetici özeti
        sl = prs.slides.add_slide(blank)
        txb(sl, "Yönetici Özeti", 0.5, 0.25, 12.3, 0.75, size=22, bold=True, color=BLUE)
        txb(sl, result.executive_summary, 0.5, 1.2, 12.3, 5.8, size=12, color=DARK)

        # Ana bulgular
        if result.key_insights:
            sl = prs.slides.add_slide(blank)
            txb(sl, "Ana Bulgular", 0.5, 0.25, 12.3, 0.75, size=22, bold=True, color=BLUE)
            y = 1.25
            for ins in result.key_insights[:7]:
                txb(sl, f"▶  {ins}", 0.7, y, 12, 0.65, size=12, color=DARK)
                y += 0.78

        # Finansal metrikler
        if result.financial_metrics:
            sl = prs.slides.add_slide(blank)
            txb(sl, "Finansal Metrikler", 0.5, 0.25, 12.3, 0.75, size=22, bold=True, color=BLUE)
            m, c = result.financial_metrics, result.financial_metrics.currency

            def v(x):
                return f"{x:,.0f} {c}" if x else "Veri yok"

            items = [
                ("Gelir",               v(m.revenue)),
                ("Net Kâr / Zarar",     v(m.net_income)),
                ("Toplam Varlık",       v(m.total_assets)),
                ("Toplam Borç",         v(m.total_liabilities)),
                ("Özkaynak",            v(m.total_equity)),
                ("İşl. Nakit Akışı",    v(m.operating_cash_flow)),
                ("Dönem",               m.period or "—"),
            ]
            y = 1.3
            for lbl, val in items:
                txb(sl, lbl, 0.7, y, 5.5, 0.55, size=12, bold=True,
                    color=RGBColor(0x61, 0x61, 0x61))
                txb(sl, val, 6.5, y, 6.2, 0.55, size=12, bold=True, color=BLUE)
                y += 0.83

        # Grafik slaytları
        for cname, cbytes in charts.items():
            if cbytes:
                sl = prs.slides.add_slide(blank)
                label = cname.replace(".png", "").replace("_", " ").title()
                txb(sl, label, 0.5, 0.2, 12.3, 0.65, size=18, bold=True, color=BLUE)
                sl.shapes.add_picture(io.BytesIO(cbytes),
                                      Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.1))

        # Öneriler
        if result.recommendations:
            sl = prs.slides.add_slide(blank)
            txb(sl, "Öneriler", 0.5, 0.25, 12.3, 0.75, size=22, bold=True, color=BLUE)
            y = 1.25
            for rec in result.recommendations[:7]:
                txb(sl, f"✓  {rec}", 0.7, y, 12, 0.65, size=12, color=GREEN)
                y += 0.82

        # Risk / Anomaliler
        if result.anomalies:
            sl = prs.slides.add_slide(blank)
            txb(sl, "Risk ve Anomaliler", 0.5, 0.25, 12.3, 0.75, size=22, bold=True, color=RED)
            y = 1.25
            for a in result.anomalies[:7]:
                txb(sl, f"⚠  {a}", 0.7, y, 12, 0.65, size=12, color=RED)
                y += 0.82

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ── Excel ─────────────────────────────────────────────────────────────

    def _excel(self, result: AnalysisResult, raw: dict) -> Optional[bytes]:
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        except ImportError:
            return None

        wb = openpyxl.Workbook()
        BLUE_FILL  = PatternFill("solid", fgColor="1565C0")
        LIGHT_FILL = PatternFill("solid", fgColor="E3F2FD")
        ALT_FILL   = PatternFill("solid", fgColor="FFFFFF")
        H_FONT = Font(bold=True, color="FFFFFF", size=11)
        T_FONT = Font(bold=True, size=13, color="1565C0")
        thin   = Side(border_style="thin", color="BDBDBD")
        brd    = Border(left=thin, right=thin, top=thin, bottom=thin)

        def hdr(ws, row, col, text):
            c = ws.cell(row=row, column=col, value=text)
            c.fill = BLUE_FILL; c.font = H_FONT
            c.alignment = Alignment(horizontal="center")

        # Sayfa 1 — Özet
        ws = wb.active; ws.title = "Özet"
        ws["A1"] = "Finansal Belge Analiz Raporu"
        ws["A1"].font = Font(bold=True, size=15, color="1565C0")
        ws["A2"] = f"Analiz Tarihi: {datetime.now().strftime('%d.%m.%Y %H:%M')} | Paket: {result.tier.value.title()} | Belge: {result.document_count}"
        ws["A2"].font = Font(size=10, color="616161")
        for cell in ["A1", "A2"]:
            ws.merge_cells(f"{cell}:E{cell[-1]}")

        ws["A4"] = "YÖNETİCİ ÖZETİ"; ws["A4"].font = T_FONT
        ws["A5"] = result.executive_summary
        ws["A5"].alignment = Alignment(wrap_text=True)
        ws.row_dimensions[5].height = 90
        ws.merge_cells("A5:E5")

        ws["A7"] = "ANA BULGULAR"; ws["A7"].font = T_FONT
        for i, ins in enumerate(result.key_insights, 8):
            ws.cell(row=i, column=1, value=f"• {ins}")
            ws.merge_cells(f"A{i}:E{i}")

        for col in "ABCDE":
            ws.column_dimensions[col].width = 28

        # Sayfa 2 — Finansal Metrikler
        if result.financial_metrics:
            ws2 = wb.create_sheet("Finansal Metrikler")
            for j, h in enumerate(["Metrik", "Değer", "Para Birimi"], 1):
                hdr(ws2, 1, j, h)
            m, c = result.financial_metrics, result.financial_metrics.currency
            rows = [
                ("Gelir",               m.revenue,              c),
                ("Gider",               m.expenses,             c),
                ("Net Kâr / Zarar",     m.net_income,           c),
                ("Toplam Varlık",       m.total_assets,         c),
                ("Toplam Borç",         m.total_liabilities,    c),
                ("Özkaynak",            m.total_equity,         c),
                ("İşletme Nakit Akışı", m.operating_cash_flow,  c),
                ("Dönem",               m.period,               "-"),
            ]
            for ri, (lbl, val, cu) in enumerate(rows, 2):
                fill = LIGHT_FILL if ri % 2 == 0 else ALT_FILL
                for ci, v in enumerate([lbl, val, cu], 1):
                    cell = ws2.cell(row=ri, column=ci, value=v)
                    cell.fill = fill; cell.border = brd
                    cell.alignment = Alignment(horizontal="left" if ci == 1 else "center")
            ws2.column_dimensions["A"].width = 28
            ws2.column_dimensions["B"].width = 20
            ws2.column_dimensions["C"].width = 14

        # Sayfa 3 — Belge Analizleri
        ws3 = wb.create_sheet("Belge Analizleri")
        for j, h in enumerate(["Belge Adı", "Tür", "Özet", "Ana Bulgular"], 1):
            hdr(ws3, 1, j, h)
        for ri, ds in enumerate(result.document_summaries, 2):
            ws3.cell(row=ri, column=1, value=ds.filename)
            ws3.cell(row=ri, column=2, value=ds.document_type)
            c3 = ws3.cell(row=ri, column=3, value=ds.brief_summary)
            c3.alignment = Alignment(wrap_text=True)
            ws3.cell(row=ri, column=4, value="\n".join(ds.key_findings))
            ws3.row_dimensions[ri].height = 55
        for col, w in zip("ABCD", [28, 14, 42, 42]):
            ws3.column_dimensions[col].width = w

        # Sayfa 4 — Rasyo Analizi (enterprise)
        ratio_data = raw.get("ratio_analysis", {})
        if ratio_data:
            ws4 = wb.create_sheet("Rasyo Analizi")
            for j, h in enumerate(["Rasyo", "Değer", "Yorum"], 1):
                hdr(ws4, 1, j, h)
            interpretations = {
                "current_ratio":    "≥1.5 sağlıklı",
                "debt_to_equity":   "<1 düşük riskli",
                "profit_margin":    "Yüksek tercih edilir",
                "return_on_equity": "Yüksek tercih edilir",
                "asset_turnover":   "≥1 verimli",
            }
            labels = {
                "current_ratio":    "Cari Oran",
                "debt_to_equity":   "Borç / Özkaynak",
                "profit_margin":    "Kâr Marjı",
                "return_on_equity": "Özkaynak Getirisi",
                "asset_turnover":   "Varlık Devir Hızı",
            }
            for ri, (key, name) in enumerate(labels.items(), 2):
                ws4.cell(row=ri, column=1, value=name)
                ws4.cell(row=ri, column=2, value=ratio_data.get(key))
                ws4.cell(row=ri, column=3, value=interpretations.get(key, ""))
            for col, w in zip("ABC", [28, 16, 24]):
                ws4.column_dimensions[col].width = w

        # Sayfa 5 — Öneriler ve Riskler (enterprise)
        if result.recommendations or result.anomalies:
            ws5 = wb.create_sheet("Öneriler ve Riskler")
            ws5["A1"] = "Öneriler"; ws5["A1"].font = Font(bold=True, color="2E7D32", size=12)
            for i, rec in enumerate(result.recommendations, 2):
                ws5.cell(row=i, column=1, value=f"✓ {rec}")

            offset = len(result.recommendations) + 4
            ws5.cell(row=offset, column=1, value="Risk ve Anomaliler")
            ws5.cell(row=offset, column=1).font = Font(bold=True, color="C62828", size=12)
            for i, a in enumerate(result.anomalies, offset + 1):
                ws5.cell(row=i, column=1, value=f"⚠ {a}")
            ws5.column_dimensions["A"].width = 70

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        return buf.read()
